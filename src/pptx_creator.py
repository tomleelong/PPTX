"""PowerPoint presentation creation utilities."""

from typing import List, Optional, Tuple, Dict, Any
from pathlib import Path
import logging

from pptx import Presentation
from pptx.slide import Slide
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

logger = logging.getLogger(__name__)


class PPTXCreator:
    """A utility class for creating PowerPoint presentations programmatically."""
    
    def __init__(self, template_path: Optional[str] = None):
        """
        Initialize a new PowerPoint presentation.
        
        Args:
            template_path: Optional path to .pptx template file for styling
        """
        if template_path and Path(template_path).exists():
            # Simple approach: just load template as-is
            # User should manually remove content slides from template file if needed
            self.presentation = Presentation(template_path)
            logger.info(f"Using template: {template_path}")
        else:
            self.presentation = Presentation()
        
    def add_title_slide(
        self, 
        title: str, 
        subtitle: Optional[str] = None
    ) -> Slide:
        """
        Add a title slide to the presentation.
        
        Args:
            title: The main title text
            subtitle: Optional subtitle text
            
        Returns:
            The created slide object
        """
        title_slide_layout = self.presentation.slide_layouts[0]
        slide = self.presentation.slides.add_slide(title_slide_layout)
        
        title_shape = slide.shapes.title
        title_shape.text = title
        
        if subtitle and len(slide.placeholders) > 1:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = subtitle
            
        logger.info(f"Added title slide: {title}")
        return slide
    
    def add_content_slide(
        self, 
        title: str, 
        content: List[str],
        layout_index: int = 1
    ) -> Slide:
        """
        Add a content slide with bullet points.
        
        Args:
            title: The slide title
            content: List of bullet point texts
            layout_index: Layout to use (default: 1 for title and content)
            
        Returns:
            The created slide object
        """
        bullet_slide_layout = self.presentation.slide_layouts[layout_index]
        slide = self.presentation.slides.add_slide(bullet_slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = title
        
        # Add content
        content_shape = slide.placeholders[1]
        text_frame = content_shape.text_frame
        text_frame.clear()
        
        for i, bullet_text in enumerate(content):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            p.text = bullet_text
            p.level = 0
            
        logger.info(f"Added content slide: {title} with {len(content)} bullet points")
        return slide
    
    def add_image_slide(
        self, 
        title: str, 
        image_path: str,
        image_width: Optional[float] = None,
        image_height: Optional[float] = None
    ) -> Slide:
        """
        Add a slide with an image.
        
        Args:
            title: The slide title
            image_path: Path to the image file
            image_width: Optional width in inches
            image_height: Optional height in inches
            
        Returns:
            The created slide object
        """
        slide_layout = self.presentation.slide_layouts[6]  # Blank layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Add title
        title_shape = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(9), Inches(1)
        )
        title_frame = title_shape.text_frame
        title_frame.text = title
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.font.size = Pt(32)
        title_paragraph.font.bold = True
        
        # Add image
        if not Path(image_path).exists():
            raise FileNotFoundError(f"Image file not found: {image_path}")
            
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(image_width) if image_width else Inches(8)
        height = Inches(image_height) if image_height else None
        
        slide.shapes.add_picture(image_path, left, top, width, height)
        
        logger.info(f"Added image slide: {title}")
        return slide
    
    def add_text_slide(
        self, 
        title: str, 
        text_content: str,
        font_size: int = 18
    ) -> Slide:
        """
        Add a slide with paragraph text content.
        
        Args:
            title: The slide title
            text_content: The main text content
            font_size: Font size for the content
            
        Returns:
            The created slide object
        """
        # Use title and content layout (layout index 1)
        slide_layout = self.presentation.slide_layouts[1]
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = title
        
        # Add content
        content_shape = slide.placeholders[1]
        text_frame = content_shape.text_frame
        text_frame.text = text_content
        
        # Format text
        for paragraph in text_frame.paragraphs:
            paragraph.font.size = Pt(font_size)
            
        logger.info(f"Added text slide: {title}")
        return slide
    
    def customize_slide_design(
        self,
        slide: Slide,
        background_color: Optional[Tuple[int, int, int]] = None,
        title_color: Optional[Tuple[int, int, int]] = None,
        title_font_size: Optional[int] = None
    ) -> None:
        """
        Customize the visual design of a slide.
        
        Args:
            slide: The slide to customize
            background_color: RGB tuple for background color
            title_color: RGB tuple for title color
            title_font_size: Title font size in points
        """
        if background_color:
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(*background_color)
            
        if slide.shapes.title and (title_color or title_font_size):
            title_shape = slide.shapes.title
            if title_color:
                title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(*title_color)
            if title_font_size:
                title_shape.text_frame.paragraphs[0].font.size = Pt(title_font_size)
                
        logger.info("Applied custom slide design")
    
    def save_presentation(self, filename: str) -> None:
        """
        Save the presentation to a file.
        
        Args:
            filename: The output filename (should end with .pptx)
        """
        if not filename.endswith('.pptx'):
            filename += '.pptx'
            
        self.presentation.save(filename)
        logger.info(f"Presentation saved as: {filename}")
    
    def get_slide_count(self) -> int:
        """Get the number of slides in the presentation."""
        return len(self.presentation.slides)
    
    def create_from_outline(self, outline: Dict[str, Any]) -> None:
        """
        Create a presentation from a structured outline.
        
        Args:
            outline: Dictionary containing presentation structure
                    Example: {
                        'title': 'Presentation Title',
                        'subtitle': 'Subtitle',
                        'slides': [
                            {'type': 'content', 'title': 'Slide 1', 'content': ['Point 1', 'Point 2']},
                            {'type': 'text', 'title': 'Slide 2', 'content': 'Paragraph text'},
                            {'type': 'image', 'title': 'Slide 3', 'image_path': 'image.png'}
                        ]
                    }
        """
        # Add title slide
        if 'title' in outline:
            self.add_title_slide(
                title=outline['title'],
                subtitle=outline.get('subtitle')
            )
        
        # Add content slides
        if 'slides' in outline:
            for slide_info in outline['slides']:
                slide_type = slide_info.get('type', 'content')
                title = slide_info.get('title', 'Untitled Slide')
                
                if slide_type == 'content' and 'content' in slide_info:
                    content = slide_info['content']
                    if isinstance(content, str):
                        content = [content]
                    self.add_content_slide(title, content)
                    
                elif slide_type == 'text' and 'content' in slide_info:
                    self.add_text_slide(title, slide_info['content'])
                    
                elif slide_type == 'image' and 'image_path' in slide_info:
                    try:
                        self.add_image_slide(
                            title, 
                            slide_info['image_path'],
                            slide_info.get('image_width'),
                            slide_info.get('image_height')
                        )
                    except FileNotFoundError as e:
                        logger.warning(f"Skipping image slide due to error: {e}")
                        # Add a text slide instead
                        self.add_text_slide(title, f"Image not found: {slide_info['image_path']}")
        
        logger.info(f"Created presentation from outline with {self.get_slide_count()} slides")